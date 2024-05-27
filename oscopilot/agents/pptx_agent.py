from oscopilot.utils import setup_config, setup_pre_run
from oscopilot.modules.base_module import BaseModule
import re
import os
import json
import base64
import time
from rich.console import Console
from rich.markdown import Markdown

from desktop_env.envs.desktop_env import DesktopEnv

console = Console()

def encode_image(image_content):
    return base64.b64encode(image_content).decode('utf-8')

def rich_print(markdown_text):
    md = Markdown(markdown_text)
    console.print(md)


def send_chat_prompts(message, llm):
    return llm.chat(message)

    
def extract_code(input_string):
    pattern = r"```(\w+)?\s*(.*?)```"  # 匹配代码块，语言名称是可选项
    matches = re.findall(pattern, input_string, re.DOTALL)

    if matches:
        language, code = matches[0]

        # 如果没有语言信息，尝试从代码内容中检测
        if not language:
            if re.search("python", code.lower()) or re.search(r"import\s+\w+", code):
                language = "python"
            elif re.search("bash", code.lower()) or re.search(r"echo", code):
                language = "bash"

        if language == 'bash':
            code = code.strip()
            code = code.replace('"', '\\"')
            code = code.replace('\n', ' && ')
            code = "pyautogui.typewrite(\"{0}\", interval=0.05)\npyautogui.press('enter')\ntime.sleep(2)".format(code)
            if re.search("sudo", code.lower()):
                code += "\npyautogui.typewrite('password', interval=0.05)\npyautogui.press('enter')\ntime.sleep(1)"
        elif language == 'python':

            # save code
            with open('tmp.py', 'w') as f:
                f.write(code)

            code = "pyautogui.typewrite(\"python main.py\", interval=0.05)\npyautogui.press('enter')\ntime.sleep(2)"
            
        else:
            raise language

        return code, language
    else:
        return None, None

from desktop_env.envs.desktop_env import DesktopEnv
class PptxAgent(BaseModule):
    info = {
            "name": "PptxAgent",
            "can do": "specializes in creating and modifying PowerPoint presentations using Python's python-pptx library. It can handle tasks involving slide creation, layout management, text and content insertion, and formatting adjustments. Also capable of detecting open PowerPoint presentations using Bash commands.",
            "can't do": "cannot handle GUI operations, cannot perform tasks outside the capabilities of the python-pptx library such as directly interacting with embedded videos and complex animations. Additionally, cannot modify LibreOffice Impress software defaults or preferences."
        }

    def __init__(self, args, task_name, env):
        super().__init__()
        self.args = args

        self.environment = env

        self.task_name = task_name

        self.reply = None

    def execute_tool(self, code, lang):
        if lang == 'python':
            file = [{
                "local_path": 'tmp.py',
                "path": '/home/user/main.py'
              }]
            self.environment.setup_controller._upload_file_setup(file)
        obs, reward, done, info = self.environment.step(code)  # node_type
        
        reply = self.environment.controller.get_terminal_output()
        # update reply
        if self.reply and reply:
            message_terminal = reply.replace(self.reply, "")
        else:
            message_terminal = reply
        self.reply = reply

        base64_image = encode_image(obs['screenshot'])
        print("message_terminal", message_terminal)
        if 'gpt' in self.llm.model_name:
            message = {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "After executing the command, the terminal output is {0}. the screenshot as below.".format(message_terminal)
                        }
                        # {
                        #     "type": "image_url",
                        #     "image_url": {
                        #         "url": f"data:image/png;base64,{base64_image}",
                        #         "detail": "high"
                        #     }
                        # }
                    ]
            }
        else:
            message = {
                    "role": "user",
                    "content": "After executing the command, the terminal output is {0}. the screenshot as below.".format(message_terminal)
            }
        return message
# When a user refers to a file opened， you can use the provided screenshot to find the filepath.
    def run(self):

        while not self.environment.controller.get_terminal_output():
            self.environment.step("pyautogui.click()")
            self.environment.step("pyautogui.hotkey('ctrl', 'alt', 't')")

        light_planner_sys_prompt = '''You are SlideAgent, an advanced programming assistant specialized in creating and modifying PowerPoint presentations. 
        Your abilities extend to manipulating slides using Python's python-pptx library.
First, write a plan. **Always recap the plan between each code block** (you have extreme short-term memory loss, so you need to recap the plan between each message block to retain it).
When you execute code, it will be executed **on the user's machine**. The user has given you **full and complete permission** to execute any code necessary to complete the task. Execute the code.
If you want to send data between programming languages, save the data to a txt or json.
You can access the internet. Run **any code** to achieve the goal, and if at first you don't succeed, try again and again.
You can install new packages.
Important: When a user refers to a file opened or a url, you can use your visual capacity or command line tools such as "lsof | grep -E '\.odp|\.pptx'" to see the path of the file being opened.
First thing need to do is pip install python-pptx.
Write messages to the user in Markdown.
In general, try to **make plans** with as few steps as possible. As for actually executing code to carry out that plan, for *stateful* languages (like python, javascript, shell, but NOT for html which starts from 0 every time) **it's critical not to try to do everything in one code block.** You should try something, print information about it, then continue from there in tiny, informed steps. You will never get it on the first try, and attempting it in one go will often lead to errors you cant see.
You are capable of **any** task.
Example Command in Bash:
```bash
pip install python-pptx && lsof | grep '.pptx'
```
Example Operation in Python:
```python
from pptx import Presentation

# Create or open the presentation
prs = Presentation()

# Adding a slide
slide = prs.slides.add_slide(prs.slide_layouts[1])  # Using the title and content layout
slide.shapes.title.text = "Hello, SlideAgent"
slide.placeholders[1].text = "This is a slide created by SlideAgent."

# Save the presentation, overwriting the original
prs.save('path_to_presentation.pptx')
print("Python Script Executed Successfully!!!")
```
Each time you only need to output the next command and wait for a reply.
'''  #  Try to use `print` or `echo` to output information needed for the subsequent tasks, or the next step might not get the required information.
        light_planner_user_prompt = '''
        User's information are as follows:
        System Version: Ubuntu
        Task Name: {task_name}
        '''.format(task_name=self.task_name)
        
        message = [
            {"role": "system", "content": light_planner_sys_prompt},
            {"role": "user", "content": light_planner_user_prompt},
        ]
 
        while True:
            print("send_chat_prompts...")
            response = send_chat_prompts(message, self.llm)
            rich_print(response)
            message.append({"role": "assistant", "content": response})

            code, lang = extract_code(response)
            if code:
                result = self.execute_tool(code, lang)
            else:
                result = ''

            if result != '':
                message.append(result)
            else:
                message.append({"role": "user", "content": "Please continue. If all tasks have been completed, reply with 'Execution Complete'. If you believe subsequent tasks cannot continue, reply with 'Execution Interrupted', including the reasons why the tasks cannot proceed, and provide the user with some possible solutions."})
            
            if 'Execution Complete' in response or 'Execution Interrupted' in response:
                break

        return response